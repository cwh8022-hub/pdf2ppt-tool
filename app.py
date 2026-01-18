import streamlit as st
from pdf2image import convert_from_bytes
from pptx import Presentation
from io import BytesIO

st.set_page_config(page_title="NotebookLM PDF 轉簡報", layout="centered")
st.title("📊 PDF 高清分頁轉 PPTX")
st.write("這是一個專為 NotebookLM 打造的工具，100% 分頁還原。")

uploaded_file = st.file_uploader("請上傳 PDF 檔案", type="pdf")

if uploaded_file:
    if st.button("🚀 開始轉換並準備下載"):
        with st.spinner('正在渲染高清分頁中...請稍候'):
            # 1. 將 PDF 轉為圖片 (Python 核心強項)
            images = convert_from_bytes(uploaded_file.read(), dpi=200)
            
            # 2. 建立 PPT
            prs = Presentation()
            for img in images:
                # 根據圖片比例動態調整投影片尺寸 (9525 是像素轉 EMU 的倍率)
                prs.slide_width = img.width * 9525
                prs.slide_height = img.height * 9525
                
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                img_io = BytesIO()
                img.save(img_io, 'JPEG', quality=95)
                img_io.seek(0)
                slide.shapes.add_picture(img_io, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # 3. 提供下載
            pptx_io = BytesIO()
            prs.save(pptx_io)
            st.download_button(
                label="📥 下載完成的 PPTX 檔案",
                data=pptx_io.getvalue(),
                file_name="Notebook_Converted.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        st.success("🎉 轉換完成！下載後直接丟進 Google Drive 即可。")
